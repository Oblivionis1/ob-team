import os
import logging
from typing import Dict, Any, List, Optional

def process_ppt(file_path: str) -> Dict[str, Any]:
    """
    处理PowerPoint文件并提取文本内容和元数据
    
    Args:
        file_path: PowerPoint文件的路径
        
    Returns:
        包含处理结果的字典，包括提取的文本、幻灯片内容和元数据
    """
    result = {
        "success": False,
        "text": "",
        "slides": [],
        "metadata": {},
        "error": None
    }
    
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            result["error"] = f"文件不存在: {file_path}"
            return result
        
        # 检查文件扩展名
        _, ext = os.path.splitext(file_path)
        if ext.lower() not in ['.ppt', '.pptx']:
            result["error"] = f"不支持的文件格式: {ext}. 仅支持PPT/PPTX文件."
            return result
        
        # 提取PPT内容
        slides, metadata = extract_content_from_ppt(file_path)
        
        # 合并所有幻灯片的文本
        all_text = []
        for slide in slides:
            all_text.append(slide["title"] if slide["title"] else "")
            all_text.append(slide["content"] if slide["content"] else "")
        
        combined_text = "\n\n".join([text for text in all_text if text])
        
        # 获取文件元数据
        file_metadata = {
            "filename": os.path.basename(file_path),
            "size_bytes": os.path.getsize(file_path),
            "slide_count": len(slides),
            "word_count": len(combined_text.split()),
            "char_count": len(combined_text)
        }
        
        # 合并PPT特定的元数据
        file_metadata.update(metadata)
        
        result["success"] = True
        result["text"] = combined_text
        result["slides"] = slides
        result["metadata"] = file_metadata
        
    except Exception as e:
        result["error"] = str(e)
        logging.error(f"处理PPT文件失败: {str(e)}", exc_info=True)
    
    return result

def extract_content_from_ppt(file_path: str) -> tuple:
    """
    从PowerPoint文件中提取内容和元数据
    
    Returns:
        tuple: (幻灯片列表, 元数据字典)
    """
    try:
        from pptx import Presentation
        
        slides = []
        metadata = {}
        
        # 打开PPT文件
        prs = Presentation(file_path)
        
        # 提取元数据
        core_properties = prs.core_properties
        metadata = {
            "title": core_properties.title or "",
            "author": core_properties.author or "",
            "subject": core_properties.subject or "",
            "keywords": core_properties.keywords or "",
            "created": core_properties.created.strftime('%Y-%m-%d %H:%M:%S') if core_properties.created else "",
            "last_modified_by": core_properties.last_modified_by or "",
            "slide_count": len(prs.slides)
        }
        
        # 提取每一张幻灯片的内容
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "title": "",
                "content": "",
                "notes": ""
            }
            
            # 提取标题
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text
            
            # 提取幻灯片中的所有文本
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text and shape.text.strip():
                        texts.append(shape.text.strip())
            
            # 合并所有文本
            slide_data["content"] = "\n".join(texts)
            
            # 提取备注
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_data["notes"] = notes_slide.notes_text_frame.text
            
            slides.append(slide_data)
        
        return slides, metadata
        
    except ImportError:
        logging.warning("未找到python-pptx库，无法处理PowerPoint文件。")
        return [], {"error": "需要安装python-pptx库以处理PowerPoint文件。"}
    except Exception as e:
        logging.error(f"处理PowerPoint文件时出错: {str(e)}", exc_info=True)
        return [], {"error": str(e)} 
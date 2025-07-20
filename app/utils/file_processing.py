import os
import PyPDF2
from pptx import Presentation
import speech_recognition as sr
from moviepy.editor import VideoFileClip
import cv2
import pytesseract

def extract_text(file_path, file_type):
    text = ""
    
    if file_type == 'pdf':
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
    
    elif file_type == 'pptx':
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    
    elif file_type in ['mp3', 'wav']:
        r = sr.Recognizer()
        with sr.AudioFile(file_path) as source:
            audio = r.record(source)
            text = r.recognize_google(audio)
    
    elif file_type in ['mp4', 'avi']:
        # 提取音频
        video = VideoFileClip(file_path)
        audio_path = "temp_audio.wav"
        video.audio.write_audiofile(audio_path)
        text += extract_text(audio_path, 'wav')
        
        # 提取视频帧中的文字
        cap = cv2.VideoCapture(file_path)
        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_interval = int(fps * 5)  # 每5秒取一帧
        
        frame_count = 0
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            if frame_count % frame_interval == 0:
                frame_text = pytesseract.image_to_string(frame)
                text += "\n" + frame_text
            frame_count += 1
        cap.release()
    
    elif file_type == 'txt':
        with open(file_path, 'r') as file:
            text = file.read()
    
    return text